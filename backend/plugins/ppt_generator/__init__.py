"""PowerPoint Generator Plugin (Phase 3 — Autonomous Document & Deck Synthesis).

Generates styled dark-theme PowerPoint (.pptx) slide decks programmatically using python-pptx,
saving LLM tokens by operating over structured JSON slide specifications.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.app.hands import registry

DARK_BG = RGBColor(15, 15, 19)
TEXT_MAIN = RGBColor(240, 240, 242)
ACCENT_BLUE = RGBColor(10, 132, 255)
TEXT_MUTED = RGBColor(142, 142, 147)


def _apply_dark_background(slide: Any) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def _ppt_create_deck(
    title: str,
    subtitle: str = "",
    slides: list[dict[str, Any]] | None = None,
    output_path: str = "presentation.pptx",
) -> dict[str, Any]:
    """Generate a native PowerPoint .pptx presentation deck."""
    clean_title = (title or "Presentation").strip()
    clean_path = Path((output_path or "presentation.pptx").strip())
    clean_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    _apply_dark_background(slide)

    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]

    title_box.text = clean_title
    for paragraph in title_box.text_frame.paragraphs:
        paragraph.font.color.rgb = ACCENT_BLUE
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT

    if subtitle:
        subtitle_box.text = subtitle
        for paragraph in subtitle_box.text_frame.paragraphs:
            paragraph.font.color.rgb = TEXT_MUTED
            paragraph.font.size = Pt(22)
            paragraph.alignment = PP_ALIGN.LEFT

    # 2. Content Slides
    content_slides = slides or []
    blank_layout = prs.slide_layouts[6]

    for slide_data in content_slides:
        s_title = str(slide_data.get("title", "Topic")).strip()
        bullets = slide_data.get("bullets", [])

        c_slide = prs.slides.add_slide(blank_layout)
        _apply_dark_background(c_slide)

        # Slide Title
        txBox = c_slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        # Slide Content / Bullets
        if bullets:
            bBox = c_slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
            btf = bBox.text_frame
            btf.word_wrap = True

            for idx, bullet in enumerate(bullets):
                bp = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
                bp.text = f"•  {str(bullet).strip()}"
                bp.font.size = Pt(20)
                bp.font.color.rgb = TEXT_MAIN
                bp.space_after = Pt(14)

    prs.save(str(clean_path))

    return {
        "ok": True,
        "title": clean_title,
        "total_slides": 1 + len(content_slides),
        "output_file": str(clean_path.resolve()),
    }


registry.register(
    {
        "name": "ppt_create_deck",
        "description": (
            "Create a native dark-theme PowerPoint presentation (.pptx). "
            "Pass title, optional subtitle, and slides array of {title, bullets}."
        ),
        "version": "1.0.0",
        "phase": 3,
        "risk_level": "auto",
        "executor": "brain",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Main title of the presentation deck",
                },
                "subtitle": {
                    "type": "string",
                    "description": "Optional subtitle for the cover slide",
                },
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                        "required": ["title"],
                    },
                    "description": "List of content slides with title and bullet points",
                },
                "output_path": {
                    "type": "string",
                    "description": "Output path for .pptx file (default 'presentation.pptx')",
                },
            },
            "required": ["title"],
        },
        "returns": {
            "type": "object",
            "properties": {
                "ok": {"type": "boolean"},
                "total_slides": {"type": "integer"},
                "output_file": {"type": "string"},
            },
        },
        "scopes": ["docs:write"],
        "tags": ["docs", "presentation", "powerpoint"],
    },
    _ppt_create_deck,
)
